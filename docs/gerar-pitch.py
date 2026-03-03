"""
Gera a apresentacao comercial (pitch deck) do FinApp em PowerPoint.
Visual moderno, dados reais, narrativa persuasiva.
14 slides — marco 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Paleta de cores ────────────────────────────────────────────────────────────
EMERALD_600 = RGBColor(0x05, 0x96, 0x69)
EMERALD_500 = RGBColor(0x10, 0xB9, 0x81)
EMERALD_400 = RGBColor(0x34, 0xD3, 0x99)
EMERALD_100 = RGBColor(0xD1, 0xFA, 0xE5)
EMERALD_50  = RGBColor(0xEC, 0xFD, 0xF5)
EMERALD_DK  = RGBColor(0x04, 0x7D, 0x57)

SLATE_900   = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800   = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700   = RGBColor(0x33, 0x41, 0x55)
SLATE_600   = RGBColor(0x47, 0x55, 0x69)
SLATE_500   = RGBColor(0x64, 0x74, 0x8B)
SLATE_400   = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300   = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200   = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100   = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50    = RGBColor(0xF8, 0xFA, 0xFC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

ROSE_600    = RGBColor(0xE1, 0x1D, 0x48)
ROSE_500    = RGBColor(0xF4, 0x3F, 0x5E)
ROSE_100    = RGBColor(0xFF, 0xE4, 0xE6)
ROSE_50     = RGBColor(0xFF, 0xF1, 0xF2)
AMBER_500   = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_100   = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_50    = RGBColor(0xFF, 0xFB, 0xEB)
BLUE_600    = RGBColor(0x25, 0x63, 0xEB)
BLUE_500    = RGBColor(0x3B, 0x82, 0xF6)
BLUE_100    = RGBColor(0xDB, 0xEA, 0xFE)
VIOLET_500  = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_50   = RGBColor(0xED, 0xE9, 0xFE)

FONT        = "Segoe UI"
W           = Inches(13.333)
H           = Inches(7.5)


# ── Helpers base ───────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, line_color=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def oval(slide, l, t, w, h, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=Pt(14), color=SLATE_700,
        bold=False, align=PP_ALIGN.LEFT, font=FONT,
        wrap=True, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set("anchor", "ctr")
    elif anchor == MSO_ANCHOR.BOTTOM:
        bodyPr.set("anchor", "b")
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    if spacing:
        p.line_spacing = spacing
    return tb


def mtxt(slide, l, t, w, h, lines, default_size=Pt(14),
         default_color=SLATE_700, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=None):
    """Multi-paragraph text box. lines = [(text, {size,color,bold,align,space_after})]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set("anchor", "ctr")
    for i, (text, s) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = s.get("font", FONT)
        p.font.size = s.get("size", default_size)
        p.font.color.rgb = s.get("color", default_color)
        p.font.bold = s.get("bold", False)
        p.alignment = s.get("align", align)
        if s.get("space_after"):
            p.space_after = s["space_after"]
        if spacing:
            p.line_spacing = spacing
    return tb


def label(slide, l, t, text, color=EMERALD_600):
    """Small uppercase section label."""
    txt(slide, l, t, Inches(5), Inches(0.4), text,
        size=Pt(10.5), color=color, bold=True)


def accent_bar(slide, l, t, w=Inches(0.7), h=Pt(4), color=EMERALD_600):
    rect(slide, l, t, w, h, color)


def stripe_left(slide, color=EMERALD_600):
    rect(slide, 0, 0, Inches(0.12), H, color)


def deco_circle(slide, l, t, size, color=SLATE_800):
    oval(slide, l, t, size, size, color)


def checklist(slide, l, t, w, items, size=Pt(13), color=EMERALD_600, text_color=SLATE_700, gap=Pt(10)):
    tb = slide.shapes.add_textbox(l, t, w, Inches(len(items) * 0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "\u2713  "
        r1.font.name = FONT
        r1.font.size = size
        r1.font.color.rgb = color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = size
        r2.font.color.rgb = text_color
        p.space_after = gap


def xlist(slide, l, t, w, items, size=Pt(12.5), color=ROSE_500, text_color=SLATE_600):
    tb = slide.shapes.add_textbox(l, t, w, Inches(len(items) * 0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "\u2717  "
        r1.font.name = FONT
        r1.font.size = size
        r1.font.color.rgb = color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = size
        r2.font.color.rgb = text_color
        p.space_after = Pt(6)


# ── Componentes reutilizaveis ──────────────────────────────────────────────────

def stat_card(slide, l, t, w, h, number, desc, accent, bg_color=WHITE, number_size=Pt(40), source=None):
    rrect(slide, l, t, w, h, bg_color, line_color=SLATE_200)
    accent_bar(slide, l + Inches(0.3), t + Inches(0.2), Inches(0.6), Pt(4), accent)
    txt(slide, l + Inches(0.3), t + Inches(0.4), w - Inches(0.5), Inches(0.75),
        number, size=number_size, color=accent, bold=True)
    txt(slide, l + Inches(0.3), t + Inches(1.25), w - Inches(0.5), h - Inches(1.7),
        desc, size=Pt(12.5), color=SLATE_600, spacing=Pt(17))
    if source:
        txt(slide, l + Inches(0.3), t + h - Inches(0.45), w - Inches(0.5), Inches(0.35),
            source, size=Pt(9), color=SLATE_400)


def feature_card(slide, l, t, w, h, icon, title, desc, accent=EMERALD_600):
    rrect(slide, l, t, w, h, WHITE, line_color=SLATE_200)
    # Icon badge
    icon_size = Inches(0.5)
    oval(slide, l + Inches(0.28), t + Inches(0.22), icon_size, icon_size, EMERALD_50)
    txt(slide, l + Inches(0.28), t + Inches(0.22), icon_size, icon_size,
        icon, size=Pt(18), color=accent, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, l + Inches(0.28), t + Inches(0.84), w - Inches(0.55), Inches(0.35),
        title, size=Pt(13.5), color=SLATE_900, bold=True)
    txt(slide, l + Inches(0.28), t + Inches(1.22), w - Inches(0.55), h - Inches(1.45),
        desc, size=Pt(11), color=SLATE_600, spacing=Pt(16))


def pill(slide, l, t, text, accent=EMERALD_600, bg_color=None):
    w = Inches(2.9)
    h = Inches(0.46)
    bg_c = bg_color or EMERALD_50
    rrect(slide, l, t, w, h, bg_c)
    txt(slide, l, t, w, h, "\u2713  " + text, size=Pt(11.5), color=accent,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — CAPA
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_900)
    stripe_left(sl)

    # Orbs decorativos
    deco_circle(sl, Inches(9.8), Inches(-1.2), Inches(5))
    deco_circle(sl, Inches(11.2), Inches(5.2), Inches(3.5))
    oval(sl, Inches(0.5), Inches(5.5), Inches(2.5), Inches(2.5), SLATE_800)

    # Logo + tagline
    txt(sl, Inches(1.2), Inches(1.6), Inches(6), Inches(1.1),
        "FinApp", size=Pt(60), color=EMERALD_500, bold=True)

    accent_bar(sl, Inches(1.2), Inches(2.8), Inches(2.5), Pt(3))

    txt(sl, Inches(1.2), Inches(3.0), Inches(8.5), Inches(0.8),
        "Gestao financeira pessoal inteligente.",
        size=Pt(28), color=WHITE)

    txt(sl, Inches(1.2), Inches(3.85), Inches(9), Inches(0.6),
        "Controle completo. Projecoes reais. Assistente com IA. Gratuito.",
        size=Pt(16), color=SLATE_400)

    # Badge "gratuito"
    rrect(sl, Inches(1.2), Inches(4.75), Inches(1.8), Inches(0.45), EMERALD_600)
    txt(sl, Inches(1.2), Inches(4.75), Inches(1.8), Inches(0.45),
        "Gratuito", size=Pt(13), color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    txt(sl, Inches(1.2), Inches(6.6), Inches(5), Inches(0.4),
        "Marco 2026", size=Pt(11), color=SLATE_500)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — O PROBLEMA (CHOQUE COM DADOS)
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    # Barra vermelha de alerta no topo
    rect(sl, 0, 0, W, Inches(0.08), ROSE_500)

    label(sl, Inches(0.8), Inches(0.35), "O CENARIO ATUAL", ROSE_500)

    mtxt(sl, Inches(0.8), Inches(0.85), Inches(11.5), Inches(1.6),
         [
             ("Metade do Brasil esta no vermelho.", {"size": Pt(40), "color": SLATE_900, "bold": True}),
         ])

    txt(sl, Inches(0.8), Inches(2.05), Inches(10.5), Inches(0.55),
        "Inadimplencia bate recorde. Dividas crescem. E a maioria nao sabe como sair.",
        size=Pt(18), color=SLATE_500)

    # 3 cards de stats
    cw, ch = Inches(3.7), Inches(3.2)
    cy = Inches(3.2)
    gx = Inches(0.4)
    sx = Inches(0.8)

    stat_card(sl, sx, cy, cw, ch,
              "81 mi", "de brasileiros inadimplentes\n— recorde historico\nem dezembro de 2025",
              ROSE_500, number_size=Pt(42), source="Serasa, Dez/2025")

    stat_card(sl, sx + cw + gx, cy, cw, ch,
              "77,5%", "das familias brasileiras\nestao endividadas\nno inicio de 2026",
              AMBER_500, number_size=Pt(42), source="CNC / Agencia Brasil")

    stat_card(sl, sx + 2*(cw + gx), cy, cw, ch,
              "48%", "dos brasileiros nao\ncontrolam o proprio\norcamento",
              BLUE_500, number_size=Pt(42), source="CNDL/SPC Brasil")

    txt(sl, Inches(0.8), Inches(6.85), Inches(11), Inches(0.35),
        "Fontes: Serasa (Dez/2025) | CNC/Agencia Brasil | CNDL/SPC Brasil",
        size=Pt(8.5), color=SLATE_400)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — POR QUE AS PESSOAS NAO CONTROLAM (4 BARREIRAS)
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_50)

    label(sl, Inches(0.8), Inches(0.4), "A RAIZ DO PROBLEMA")

    txt(sl, Inches(0.8), Inches(0.85), Inches(10), Inches(0.75),
        "Por que as pessoas nao controlam suas financas?",
        size=Pt(32), color=SLATE_900, bold=True)

    barriers = [
        ("\u2716", "Nao sabem por onde comecar",
         "55% dos brasileiros entendem pouco ou nada sobre financas pessoais. "
         "19% dos jovens adultos dizem nunca ter aprendido.",
         ROSE_500, ROSE_50),

        ("\u23F3", "Falta de disciplina e continuidade",
         "Registrar gastos manualmente e tedioso. A maioria desiste em semanas. "
         "36% dos que tentam usam caderno de papel ou planilha — e abandonam.",
         AMBER_500, AMBER_50),

        ("\u2699", "Ferramentas complexas e fragmentadas",
         "Gastos em um app, investimentos em outro, projecoes numa planilha. "
         "Nenhuma ferramenta centraliza tudo com clareza e sem curva de aprendizado.",
         BLUE_500, BLUE_100),

        ("\u20AC", "Funcionalidades essenciais atras de paywall",
         "Organizze cobra R$ 35/mes so para controle basico. "
         "Mobills exige plano PRO para IA. O usuario paga mas nao engaja.",
         VIOLET_500, VIOLET_50),
    ]

    cw, ch = Inches(5.6), Inches(1.85)
    gx, gy = Inches(0.5), Inches(0.35)
    sx, sy = Inches(0.8), Inches(2.0)

    for i, (icon, title, desc, accent, bg_c) in enumerate(barriers):
        col = i % 2
        row = i // 2
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)

        rrect(sl, x, y, cw, ch, bg_c, line_color=SLATE_200)
        txt(sl, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
            icon, size=Pt(18), color=accent, bold=True)
        txt(sl, x + Inches(0.9), y + Inches(0.18), cw - Inches(1.1), Inches(0.42),
            title, size=Pt(15), color=SLATE_900, bold=True)
        txt(sl, x + Inches(0.9), y + Inches(0.65), cw - Inches(1.1), ch - Inches(0.9),
            desc, size=Pt(11.5), color=SLATE_600, spacing=Pt(17))

    txt(sl, Inches(0.8), Inches(6.88), Inches(10), Inches(0.35),
        "Fontes: CNDL/SPC Brasil | FEBRABAN | Pesquisa de Educacao Financeira (ENEF)",
        size=Pt(8.5), color=SLATE_400)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 4 — CONCORRENTES FALHAM
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    label(sl, Inches(0.8), Inches(0.35), "O MERCADO HOJE")

    txt(sl, Inches(0.8), Inches(0.78), Inches(11.5), Inches(0.65),
        "Nenhuma ferramenta resolve o problema completo.",
        size=Pt(30), color=SLATE_900, bold=True)

    txt(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
        "Cada uma ataca um pedaco — o usuario fica sem visao integrada.",
        size=Pt(15), color=SLATE_500)

    competitors = [
        ("Mobills",          "Gratuito\nlimitado",
         AMBER_500, AMBER_50,
         ["IA so no plano PRO (R$ 25/mes)", "Investimentos em app separado",
          "Sem simuladores financeiros", "Sem Metas e Dividas integrados"]),
        ("Organizze",        "R$ 35/mes",
         ROSE_500, ROSE_50,
         ["Sem assistente IA", "Sem investimentos", "Sem importacao PDF",
          "Sem fluxo de caixa projetado"]),
        ("GuiaBolso",        "Encerrado\nem 2022",
         SLATE_500, SLATE_100,
         ["Descontinuado", "Usuarios perderam dados", "Confianca destruida",
          "Sem suporte ou atualizacoes"]),
        ("Minhas Economias", "Gratuito\n(muito limitado)",
         BLUE_500, BLUE_100,
         ["Interface desatualizada", "Sem IA ou projecoes reais",
          "Instabilidade frequente", "Sem metas nem simuladores"]),
    ]

    cw, ch = Inches(2.85), Inches(4.3)
    sx, sy = Inches(0.8), Inches(2.1)
    gx = Inches(0.3)

    for i, (name, price, accent, bg_c, issues) in enumerate(competitors):
        x = sx + i * (cw + gx)
        rrect(sl, x, sy, cw, ch, bg_c, line_color=SLATE_200)

        # Nome
        txt(sl, x + Inches(0.25), sy + Inches(0.2), cw - Inches(0.4), Inches(0.55),
            name, size=Pt(17), color=SLATE_900, bold=True)

        # Preco
        rrect(sl, x + Inches(0.25), sy + Inches(0.82), cw - Inches(0.5), Inches(0.48),
              accent)
        txt(sl, x + Inches(0.25), sy + Inches(0.82), cw - Inches(0.5), Inches(0.48),
            price, size=Pt(11), color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Issues
        xlist(sl, x + Inches(0.25), sy + Inches(1.5), cw - Inches(0.35),
              issues, size=Pt(11), color=ROSE_500, text_color=SLATE_700)

    # Conclusao
    rrect(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.6), SLATE_900)
    txt(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.6),
        "  Nenhum oferece: IA integrada + Metas + Dividas + Simuladores + Importacao PDF + Plano gratuito completo",
        size=Pt(13), color=EMERALD_400, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 5 — COMPARATIVO DETALHADO
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    label(sl, Inches(0.8), Inches(0.3), "COMPARATIVO")

    txt(sl, Inches(0.8), Inches(0.68), Inches(11), Inches(0.55),
        "FinApp vs. concorrentes: funcionalidade por funcionalidade.",
        size=Pt(26), color=SLATE_900, bold=True)

    # Cabecalho da tabela
    hy = Inches(1.45)
    lx = Inches(0.7)
    col_feat = Inches(3.3)
    col_comp = Inches(4.8)
    col_fa   = Inches(3.3)
    total_w  = col_feat + col_comp + col_fa

    rect(sl, lx, hy, total_w, Inches(0.48), SLATE_900)
    txt(sl, lx + Inches(0.15), hy, col_feat, Inches(0.48),
        "Funcionalidade", size=Pt(11.5), color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(sl, lx + col_feat, hy, col_comp, Inches(0.48),
        "Concorrentes", size=Pt(11.5), color=SLATE_400, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(sl, lx + col_feat + col_comp, hy, col_fa, Inches(0.48),
        "FinApp", size=Pt(11.5), color=EMERALD_400, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("Plano gratuito completo",       "Limitado ou pago (R$ 8–35/mes)",           "\u2713  Sim, sem restricoes"),
        ("Assistente com IA",             "So Mobills (WhatsApp, plano PRO)",          "\u2713  Chat integrado (Gemini)"),
        ("Importacao OFX / CSV / PDF",    "Parcial (so OFX, so plano pago)",           "\u2713  3 formatos + IA no PDF"),
        ("KPIs e alertas inteligentes",   "Ausente ou basico",                         "\u2713  5 KPIs + insights proativos"),
        ("Tetos de orcamento",            "Ausente ou rigido",                         "\u2713  Por categoria + alertas"),
        ("Metas financeiras",             "Ausente",                                   "\u2713  CRUD + progresso + prazo"),
        ("Gestao de dividas",             "Ausente",                                   "\u2713  CRUD + simulador de quitacao"),
        ("Simuladores educacionais",      "Ausente",                                   "\u2713  4 simuladores interativos"),
        ("Fluxo de caixa projetado",      "Basico ou ausente",                         "\u2713  Diario + Previsto"),
        ("Investimentos integrados",      "Ausente ou em app separado",                "\u2713  CRUD + evolucao + retorno real"),
        ("Historico de KPIs mensal",      "Ausente",                                   "\u2713  Graficos + fechamento guiado"),
        ("Deteccao automatica de padroes","Ausente",                                   "\u2713  Sugestao de recorrentes"),
    ]

    rh   = Inches(0.44)
    ry   = hy + Inches(0.52)
    fa_x = lx + col_feat + col_comp

    for i, (feat, comp, finapp) in enumerate(rows):
        y = ry + i * rh
        # Fundo alternado
        row_bg = SLATE_50 if i % 2 == 0 else WHITE
        rect(sl, lx, y, col_feat + col_comp, rh, row_bg)

        # FinApp coluna sempre verde
        fa_bg = EMERALD_50 if i % 2 == 0 else RGBColor(0xF0, 0xFD, 0xF4)
        rect(sl, fa_x, y, col_fa, rh, fa_bg)

        txt(sl, lx + Inches(0.15), y, col_feat - Inches(0.15), rh,
            feat, size=Pt(11), color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        txt(sl, lx + col_feat, y, col_comp, rh,
            comp, size=Pt(10.5), color=SLATE_500,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(sl, fa_x, y, col_fa, rh,
            finapp, size=Pt(11), color=EMERALD_600,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 6 — A SOLUCAO (HERO)
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_900)
    stripe_left(sl)

    deco_circle(sl, Inches(10.2), Inches(0.5), Inches(5))
    deco_circle(sl, Inches(0.5), Inches(5.8), Inches(2.5))

    label(sl, Inches(1.0), Inches(0.45), "A SOLUCAO", EMERALD_400)

    txt(sl, Inches(1.0), Inches(1.1), Inches(9), Inches(1.6),
        "Tudo o que voce precisa.\nEm um unico lugar.",
        size=Pt(44), color=WHITE, bold=True, spacing=Pt(56))

    txt(sl, Inches(1.0), Inches(3.1), Inches(9), Inches(0.9),
        "Contas, transacoes, investimentos, metas, dividas, simuladores\n"
        "e um assistente com inteligencia artificial — integrados numa\n"
        "plataforma moderna. Gratuita. Sem paywall.",
        size=Pt(16.5), color=SLATE_400, spacing=Pt(24))

    # Pills 4x3
    pills_items = [
        "Controle completo",    "Importacao OFX/CSV/PDF", "Transacoes planejadas",
        "Orcamento com tetos",  "Metas financeiras",      "Gestao de dividas",
        "Fluxo de caixa",       "Investimentos + IPCA",   "Historico de KPIs",
        "Assistente IA",        "4 Simuladores",          "Deteccao de padroes",
    ]

    pill_y0 = Inches(4.5)
    pill_w  = Inches(2.9)
    pill_h  = Inches(0.44)
    gap_x   = Inches(0.42)
    gap_y   = Inches(0.2)
    pill_x0 = Inches(1.0)

    for i, p_text in enumerate(pills_items):
        col = i % 4
        row = i // 4
        x = pill_x0 + col * (pill_w + gap_x)
        y = pill_y0 + row * (pill_h + gap_y)
        rrect(sl, x, y, pill_w, pill_h, SLATE_800)
        txt(sl, x, y, pill_w, pill_h,
            "\u2713  " + p_text, size=Pt(11), color=EMERALD_400,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 7 — CONTROLE FINANCEIRO + ORCAMENTO
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_50)

    label(sl, Inches(0.8), Inches(0.38), "MODULO 1 — CONTROLE E ORCAMENTO")

    txt(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.6),
        "Seu dinheiro, dia a dia — com clareza e sem esforco.",
        size=Pt(28), color=SLATE_900, bold=True)

    # 6 feature cards em 3x2
    feats = [
        ("$", "Contas e Transacoes",
         "Multiplas contas (banco, cartao, carteira). Saldo atualizado automaticamente. Historico paginado com busca e filtros avancados."),
        ("\u2191\u2193", "Importacao Inteligente",
         "Importe OFX, CSV ou faturas PDF. Gemini extrai os dados automaticamente. Deteccao de duplicatas e auto-categorizacao por regras."),
        ("\u2637", "Orcamento por Categoria",
         "Tetos mensais por categoria de despesa. Badges de alerta (Atencao / Estourado) no Dashboard. Insights proativos baseados nos desvios."),
        ("\u21BB", "Transacoes Planejadas",
         "Recorrentes sem prazo, pontuais ou com periodo. Deteccao automatica de padroes e sugestao de criar recorrentes."),
        ("\u25A6", "Dashboard com KPIs",
         "5 indicadores-chave: taxa de poupanca, runway financeiro, reserva de emergencia, desvio orcamentario e percentual de gasto fixo."),
        ("\u2261", "Fluxo de Caixa",
         "Fluxo Diario (dia a dia com saldo acumulado) e Fluxo Previsto (projecao de 4 meses). Ideal para evitar surpresas no fim do mes."),
    ]

    cw, ch = Inches(3.7), Inches(2.5)
    gx, gy = Inches(0.35), Inches(0.28)
    sx, sy = Inches(0.8), Inches(1.65)

    for i, (icon, title, desc) in enumerate(feats):
        col = i % 3
        row = i // 3
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)
        feature_card(sl, x, y, cw, ch, icon, title, desc)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 8 — METAS, DIVIDAS E SIMULADORES
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    label(sl, Inches(0.8), Inches(0.38), "MODULO 2 — PLANEJAMENTO FINANCEIRO")

    txt(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.6),
        "De onde voce esta para onde quer chegar.",
        size=Pt(28), color=SLATE_900, bold=True)

    # Coluna esquerda — Metas
    rect(sl, Inches(0.7), Inches(1.65), Inches(3.9), Inches(5.5), EMERALD_50)
    accent_bar(sl, Inches(0.7), Inches(1.65), Inches(3.9), Pt(5), EMERALD_600)
    txt(sl, Inches(1.0), Inches(1.9), Inches(3.3), Inches(0.45),
        "Metas Financeiras", size=Pt(18), color=EMERALD_600, bold=True)
    txt(sl, Inches(1.0), Inches(2.5), Inches(3.3), Inches(0.55),
        "Crie objetivos com prazo e valor alvo.\nAcompanhe o progresso com barra visual.",
        size=Pt(12.5), color=SLATE_600, spacing=Pt(17))
    checklist(sl, Inches(1.0), Inches(3.35), Inches(3.3),
              ["Vinculada a conta real",
               "Progresso automatico pelo saldo",
               "Aporte mensal necessario",
               "Widget no Dashboard",
               "Exemplos: viagem, carro, IF"],
              size=Pt(12), gap=Pt(8))

    # Coluna centro — Dividas
    rect(sl, Inches(4.8), Inches(1.65), Inches(3.9), Inches(5.5), ROSE_50)
    accent_bar(sl, Inches(4.8), Inches(1.65), Inches(3.9), Pt(5), ROSE_500)
    txt(sl, Inches(5.1), Inches(1.9), Inches(3.3), Inches(0.45),
        "Gestao de Dividas", size=Pt(18), color=ROSE_500, bold=True)
    txt(sl, Inches(5.1), Inches(2.5), Inches(3.3), Inches(0.55),
        "Centralize emprestimos, financiamentos\ne cartoes parcelados em um lugar so.",
        size=Pt(12.5), color=SLATE_600, spacing=Pt(17))
    checklist(sl, Inches(5.1), Inches(3.35), Inches(3.3),
              ["Saldo devedor + parcelas restantes",
               "Simulador de pagamento extra",
               "Calculo de juros totais",
               "Widget no Dashboard",
               "Priorize a divida certa"],
              size=Pt(12), color=ROSE_500, gap=Pt(8))

    # Coluna direita — Simuladores
    rect(sl, Inches(8.9), Inches(1.65), Inches(3.9), Inches(5.5), VIOLET_50)
    accent_bar(sl, Inches(8.9), Inches(1.65), Inches(3.9), Pt(5), VIOLET_500)
    txt(sl, Inches(9.2), Inches(1.9), Inches(3.3), Inches(0.45),
        "Simuladores Educacionais", size=Pt(18), color=VIOLET_500, bold=True)
    txt(sl, Inches(9.2), Inches(2.5), Inches(3.3), Inches(0.55),
        "4 calculadoras interativas para\ndecisoes financeiras mais inteligentes.",
        size=Pt(12.5), color=SLATE_600, spacing=Pt(17))
    checklist(sl, Inches(9.2), Inches(3.35), Inches(3.3),
              ["Juros Compostos",
               "Impacto da Inflacao",
               "Custo de Oportunidade",
               "Independencia Financeira (FIRE)",
               "Calculos em tempo real"],
              size=Pt(12), color=VIOLET_500, gap=Pt(8))

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 9 — INVESTIMENTOS E HISTORICO
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_50)

    label(sl, Inches(0.8), Inches(0.38), "MODULO 3 — INVESTIMENTOS E ANALISE")

    txt(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.6),
        "Patrimonio e evolucao financeira em um so painel.",
        size=Pt(28), color=SLATE_900, bold=True)

    # Cards lado a lado
    # Investimentos (esquerda)
    rrect(sl, Inches(0.8), Inches(1.65), Inches(5.8), Inches(5.5), WHITE, line_color=SLATE_200)
    accent_bar(sl, Inches(1.1), Inches(2.0), Inches(0.7), Pt(4), EMERALD_600)
    txt(sl, Inches(1.1), Inches(2.18), Inches(5), Inches(0.5),
        "Carteira de Investimentos", size=Pt(20), color=SLATE_900, bold=True)
    txt(sl, Inches(1.1), Inches(2.78), Inches(5), Inches(0.7),
        "CDB, Tesouro Direto, Acoes, Cripto, Fundos — todos em um unico lugar. "
        "Aportes, resgates e atualizacoes de saldo. Quadro de evolucao mensal.",
        size=Pt(12.5), color=SLATE_600, spacing=Pt(18))
    checklist(sl, Inches(1.1), Inches(3.75), Inches(5.1),
              ["Agrupamento por tipo de produto",
               "Historico de lancamentos",
               "Evolucao mensal (quadro visual)",
               "Retorno real descontando inflacao (IPCA 12 meses)",
               "Widget com saldo total no Dashboard"],
              size=Pt(12.5), gap=Pt(8))

    # Historico de KPIs (direita)
    rrect(sl, Inches(6.9), Inches(1.65), Inches(5.8), Inches(5.5), WHITE, line_color=SLATE_200)
    accent_bar(sl, Inches(7.2), Inches(2.0), Inches(0.7), Pt(4), BLUE_500)
    txt(sl, Inches(7.2), Inches(2.18), Inches(5), Inches(0.5),
        "Historico de KPIs Mensal", size=Pt(20), color=SLATE_900, bold=True)
    txt(sl, Inches(7.2), Inches(2.78), Inches(5), Inches(0.7),
        "Acompanhe a evolucao financeira mes a mes. Cada fechamento mensal "
        "gera um snapshot automatico de todos os indicadores-chave.",
        size=Pt(12.5), color=SLATE_600, spacing=Pt(18))
    checklist(sl, Inches(7.2), Inches(3.75), Inches(5.1),
              ["Grafico Receitas vs Despesas vs Saldo",
               "Grafico de saude financeira (taxa de poupanca)",
               "Tabela comparativa mes a mes",
               "Fechamento mensal guiado com sugestoes",
               "Base para decisoes de medio prazo"],
              size=Pt(12.5), color=BLUE_500, gap=Pt(8))

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 10 — ASSISTENTE IA (DIFERENCIAL)
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    # Painel escuro esquerda
    rect(sl, 0, 0, Inches(6.5), H, SLATE_900)
    stripe_left(sl)

    deco_circle(sl, Inches(3.5), Inches(5.0), Inches(3.5), SLATE_800)

    label(sl, Inches(0.9), Inches(0.65), "O DIFERENCIAL", EMERALD_400)

    txt(sl, Inches(0.9), Inches(1.4), Inches(5.2), Inches(1.6),
        "Um assistente\nque conhece\nsuas financas.",
        size=Pt(38), color=WHITE, bold=True, spacing=Pt(50))

    txt(sl, Inches(0.9), Inches(3.6), Inches(5.0), Inches(1.1),
        "O Assistente analisa suas contas, transacoes, recorrentes, "
        "investimentos e projecoes em tempo real — e responde com "
        "diagnosticos personalizados, nao respostas genericas.",
        size=Pt(13.5), color=SLATE_400, spacing=Pt(21))

    # Badge powered by
    rrect(sl, Inches(0.9), Inches(5.0), Inches(2.8), Inches(0.42), SLATE_800)
    txt(sl, Inches(0.9), Inches(5.0), Inches(2.8), Inches(0.42),
        "Powered by Gemini 2.5 Flash", size=Pt(11),
        color=SLATE_400, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Painel direito — perguntas
    txt(sl, Inches(6.9), Inches(0.65), Inches(5.5), Inches(0.4),
        "O usuario pergunta:", size=Pt(13.5), color=SLATE_500, bold=True)

    questions = [
        "\"Como esta minha saude financeira?\"",
        "\"Minhas despesas estao controladas?\"",
        "\"Quais categorias estouraram o teto?\"",
        "\"Vale a pena quitar minha divida agora?\"",
        "\"Estou no caminho certo para minha meta?\"",
        "\"Como posso economizar mais este mes?\"",
    ]

    for i, q in enumerate(questions):
        y = Inches(1.3) + i * Inches(0.88)
        rrect(sl, Inches(6.9), y, Inches(5.7), Inches(0.68), EMERALD_50)
        txt(sl, Inches(7.15), y, Inches(5.2), Inches(0.68),
            q, size=Pt(13.5), color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

    rrect(sl, Inches(6.9), Inches(6.7), Inches(5.7), Inches(0.5), SLATE_100)
    txt(sl, Inches(7.1), Inches(6.7), Inches(5.3), Inches(0.5),
        "\u21BB  Mantem contexto conversacional — perguntas de acompanhamento naturais",
        size=Pt(11), color=SLATE_600, anchor=MSO_ANCHOR.MIDDLE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 11 — A EVIDENCIA
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_50)

    label(sl, Inches(0.8), Inches(0.38), "A EVIDENCIA")

    txt(sl, Inches(0.8), Inches(0.82), Inches(11), Inches(0.65),
        "Controlar financas funciona. Os numeros comprovam.",
        size=Pt(32), color=SLATE_900, bold=True)

    evidence = [
        ("88%",   "dos usuarios de\napps financeiros\nconsideram a\nferramenta muito\nou extremamente util",
         EMERALD_600, "Academy Bank Research"),
        ("2,5x",  "mais chance de\npoupar o suficiente\npara a aposentadoria\nquem usa planejamento\nfinanceiro",
         BLUE_500, "Ramsey Solutions"),
        ("+59%",  "crescimento em\ninstalacoes de apps\nfinanceiros na\nAmerica Latina\nem 2025",
         EMERALD_600, "TI Inside / Adjust"),
        ("90%",   "dos brasileiros\nadmitem precisar\nde mais educacao\nfinanceira",
         AMBER_500, "Funpresp-Jud"),
    ]

    cw, ch = Inches(2.9), Inches(3.8)
    sx, sy = Inches(0.8), Inches(2.3)
    gx = Inches(0.3)

    for i, (num, lbl, accent, source) in enumerate(evidence):
        x = sx + i * (cw + gx)
        stat_card(sl, x, sy, cw, ch, num, lbl, accent, number_size=Pt(42), source=source)

    # Bottom insight bar
    rrect(sl, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.68), SLATE_900)
    txt(sl, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.68),
        "  O mercado de fintech pessoal cresce 2 digitos ao ano na America Latina. "
        "O usuario quer controle — faltava a ferramenta certa.",
        size=Pt(13), color=EMERALD_400, anchor=MSO_ANCHOR.MIDDLE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 12 — OPORTUNIDADE DE MERCADO
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, WHITE)

    label(sl, Inches(0.8), Inches(0.38), "OPORTUNIDADE")

    txt(sl, Inches(0.8), Inches(0.82), Inches(11), Inches(0.65),
        "Um mercado enorme com lacunas claras.",
        size=Pt(32), color=SLATE_900, bold=True)

    # Numeros de mercado (esquerda)
    market = [
        ("USD 21,4 bi",  "Mercado global de apps\nfinanceiros pessoais em 2025", EMERALD_600),
        ("42 milhoes",   "Brasileiros ja conectados\nao Open Finance",             BLUE_500),
        ("44%",          "dos bancarizados gerem\nfinancas so pelo celular",        AMBER_500),
        ("213 milhoes",  "de brasileiros — mercado\nenderecavel domestico",         VIOLET_500),
    ]

    for i, (num, lbl, accent) in enumerate(market):
        y = Inches(1.85) + i * Inches(1.35)
        rect(sl, Inches(0.8), y, Pt(5), Inches(1.0), accent)
        txt(sl, Inches(1.1), y, Inches(4.5), Inches(0.62),
            num, size=Pt(28), color=accent, bold=True)
        txt(sl, Inches(1.1), y + Inches(0.6), Inches(4.5), Inches(0.6),
            lbl, size=Pt(12.5), color=SLATE_600, spacing=Pt(17))

    # Painel direito — lacunas que FinApp preenche
    rrect(sl, Inches(6.5), Inches(1.65), Inches(6.0), Inches(5.6), SLATE_900)
    txt(sl, Inches(7.0), Inches(2.05), Inches(5.0), Inches(0.5),
        "FinApp preenche as lacunas:", size=Pt(16),
        color=EMERALD_400, bold=True)

    gaps = [
        "Gratuito e completo — sem paywall em funcionalidades essenciais",
        "IA integrada no app — nao em canal separado (WhatsApp)",
        "Importacao OFX, CSV e PDF com IA — 3 formatos num so lugar",
        "Metas e Dividas integrados — ausentes em todos os concorrentes",
        "4 Simuladores educacionais — nao encontrado em nenhum rival",
        "KPIs, tetos e alertas — orcamento ativo, nao passivo",
        "Investimentos na mesma plataforma — sem app adicional",
        "Fechamento mensal e historico — memoria financeira real",
    ]

    checklist(sl, Inches(7.0), Inches(2.75), Inches(5.0),
              gaps, size=Pt(12), color=EMERALD_500, text_color=WHITE, gap=Pt(6))

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 13 — POR QUE FINAPP (3 PILARES)
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, EMERALD_600)

    # Decorativos
    deco_circle(sl, Inches(10.5), Inches(-0.5), Inches(3), EMERALD_DK)
    deco_circle(sl, Inches(-0.5), Inches(5.5), Inches(2.5), EMERALD_DK)

    label(sl, Inches(0.8), Inches(0.45), "POR QUE FINAPP?", EMERALD_100)

    txt(sl, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
        "Tres razoes que fazem a diferenca.",
        size=Pt(36), color=WHITE, bold=True)

    pillars = [
        ("Completo",
         "12 modulos integrados:\nContas, Transacoes, Recorrentes,\nMetas, Dividas, Investimentos,\n"
         "Fluxo, Historico, Simuladores,\nAssistente IA, Importacao e Dashboard.\n\n"
         "Uma plataforma. Sem fragmentacao."),
        ("Inteligente",
         "Assistente IA com seus dados reais.\n5 KPIs com alertas automaticos.\n"
         "Tetos de orcamento por categoria.\nDeteccao de recorrencias.\n"
         "Fechamento mensal guiado.\nRetorno real descontando inflacao."),
        ("Acessivel",
         "Interface moderna e intuitiva.\nDark mode nativo.\n"
         "Web + celular — sem instalar nada.\n\n"
         "Gratuito.\nSem paywall. Sem restricoes.\nFuncionalidades que custam\nR$ 35/mes em outros apps."),
    ]

    cw, ch = Inches(3.7), Inches(4.0)
    gx     = Inches(0.45)
    sx, sy = Inches(0.8), Inches(2.5)

    for i, (title, desc) in enumerate(pillars):
        x = sx + i * (cw + gx)
        rrect(sl, x, sy, cw, ch, EMERALD_DK)

        # Numero do pilar
        oval(sl, x + Inches(0.3), sy + Inches(0.25),
             Inches(0.45), Inches(0.45), EMERALD_400)
        txt(sl, x + Inches(0.3), sy + Inches(0.25), Inches(0.45), Inches(0.45),
            str(i + 1), size=Pt(15), color=SLATE_900,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        txt(sl, x + Inches(0.88), sy + Inches(0.22), cw - Inches(1.1), Inches(0.5),
            title, size=Pt(22), color=WHITE, bold=True)

        accent_bar(sl, x + Inches(0.3), sy + Inches(0.88), Inches(0.9), Pt(3), EMERALD_400)

        txt(sl, x + Inches(0.3), sy + Inches(1.1), cw - Inches(0.55), ch - Inches(1.4),
            desc, size=Pt(12.5), color=EMERALD_100, spacing=Pt(19))

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 14 — ENCERRAMENTO / CTA
    # ──────────────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    bg(sl, SLATE_900)
    stripe_left(sl)

    # Orbs
    deco_circle(sl, Inches(8.5), Inches(2.5), Inches(6), SLATE_800)
    deco_circle(sl, Inches(0.5), Inches(-0.5), Inches(2.5), SLATE_800)

    txt(sl, Inches(1.0), Inches(1.5), Inches(9), Inches(1.6),
        "Assuma o controle\ndas suas financas.",
        size=Pt(48), color=WHITE, bold=True, spacing=Pt(62))

    txt(sl, Inches(1.0), Inches(3.5), Inches(7.5), Inches(0.55),
        "Comece hoje. E gratuito. Sem paywall. Sem complicacao.",
        size=Pt(20), color=EMERALD_400)

    # 3 bullet points de fechamento
    closing = [
        "12 modulos integrados — tudo o que voce precisa, em um lugar",
        "Assistente IA com seus dados reais — respostas de verdade",
        "Gratuito — funcionalidades que custam R$ 35/mes nos concorrentes",
    ]

    checklist(sl, Inches(1.0), Inches(4.3), Inches(8),
              closing, size=Pt(15), color=EMERALD_500, text_color=SLATE_300, gap=Pt(10))

    # Botao CTA
    rrect(sl, Inches(1.0), Inches(5.9), Inches(3.8), Inches(0.78), EMERALD_600)
    txt(sl, Inches(1.0), Inches(5.9), Inches(3.8), Inches(0.78),
        "Conheca o FinApp  \u2192", size=Pt(19),
        color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Rodape
    accent_bar(sl, Inches(1.0), Inches(7.05), Inches(1.5), Pt(2))
    txt(sl, Inches(1.0), Inches(7.15), Inches(8), Inches(0.3),
        "FinApp  |  Gestao Financeira Pessoal  |  Marco 2026  |  finapp-kohl.vercel.app",
        size=Pt(10.5), color=SLATE_500)

    # ── Salvar ────────────────────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FinApp - Pitch Deck.pptx")
    prs.save(out)
    print("Pitch deck gerado: " + out)


if __name__ == "__main__":
    build()
